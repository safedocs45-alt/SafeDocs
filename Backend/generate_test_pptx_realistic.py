# generate_test_pptx_realistic.py
from pptx import Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout

tf = slide.shapes.add_textbox(1000000, 1000000, 7000000, 2000000).text_frame
tf.text = "SafeDocs Realistic Demo - PPTX (INERT JS STUBS)"
tf.add_paragraph().text = "Marker: X-SAFEDOCS-MARKER: MAL_TEST_JS_STUB"
tf.add_paragraph().text = "--- BEGIN PSEUDO-JS-STUB ---"
tf.add_paragraph().text = ("// PSEUDO-JS: inert stub (text only)\n"
                           "// function exploit() { /* harmless demo-only text */ }")
tf.add_paragraph().text = "--- END PSEUDO-JS-STUB ---"
tf.add_paragraph().text = "Embedded object placeholder: OLE_PAYLOAD_STUB"

prs.save('safedocs_realistic_pptx.pptx')
print("Created safedocs_realistic_pptx.pptx")
